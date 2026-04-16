"""
Fix legibility and overlap issues in Final_Presentation_kmylavarapu3.pptx.

- Removes empty text boxes (often stacked on the same coordinates).
- Sets explicit light text on all runs (avoids theme black on navy backgrounds).
- Repositions and expands content on the Systematic Model Ladder slide (slide 10).
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
PPTX_PATH = ROOT / "Final_Presentation_kmylavarapu3.pptx"

BODY = RGBColor(236, 242, 255)
TITLE = RGBColor(255, 255, 255)


def title_shape_index(slide) -> int | None:
    """Topmost non-empty TEXT_BOX is treated as the slide title."""
    candidates: list[tuple[int, int]] = []
    for idx, sh in enumerate(slide.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if not hasattr(sh, "text_frame"):
            continue
        t = sh.text_frame.text.strip()
        if not t:
            continue
        candidates.append((int(sh.top), idx))
    if not candidates:
        return None
    return min(candidates, key=lambda x: x[0])[1]


def remove_empty_textboxes(slide) -> None:
    for sh in list(slide.shapes):
        if sh.shape_type != MSO_SHAPE_TYPE.TEXT_BOX:
            continue
        if not hasattr(sh, "text_frame"):
            continue
        if sh.text_frame.text and sh.text_frame.text.strip():
            continue
        el = sh._element  # noqa: SLF001
        el.getparent().remove(el)


def apply_text_colors(slide, title_idx: int | None) -> None:
    for idx, sh in enumerate(slide.shapes):
        if not hasattr(sh, "text_frame"):
            continue
        for p in sh.text_frame.paragraphs:
            for run in p.runs:
                if not run.text:
                    continue
                if title_idx is not None and idx == title_idx:
                    run.font.color.rgb = TITLE
                    if run.font.size is None:
                        run.font.size = Pt(22)
                else:
                    run.font.color.rgb = BODY
                    if run.font.size is None:
                        run.font.size = Pt(12)


def fix_systematic_ladder_slide(slide) -> None:
    """Slide 10 — two-column layout with tall enough boxes for wrapped lines."""
    for sh in slide.shapes:
        if not hasattr(sh, "text_frame"):
            continue
        txt = sh.text_frame.text
        if "Validated Model Ladder" in txt:
            sh.left = int(Inches(0.55))
            sh.top = int(Inches(1.35))
            sh.width = int(Inches(4.35))
            sh.height = int(Inches(1.55))
        elif txt.strip().startswith("Feature Sets"):
            sh.left = int(Inches(5.05))
            sh.top = int(Inches(1.35))
            sh.width = int(Inches(4.35))
            sh.height = int(Inches(1.15))
        elif "Office-Hours Direction" in txt:
            sh.left = int(Inches(0.55))
            sh.top = int(Inches(3.05))
            sh.width = int(Inches(4.35))
            sh.height = int(Inches(1.25))
        elif txt.strip().startswith("Validation"):
            sh.left = int(Inches(5.05))
            sh.top = int(Inches(3.05))
            sh.width = int(Inches(4.35))
            sh.height = int(Inches(1.05))
        elif "Takeaway" in txt and "Halftime state" in txt:
            sh.left = int(Inches(0.65))
            sh.top = int(Inches(4.55))
            sh.width = int(Inches(8.7))
            sh.height = int(Inches(1.05))


def polish_presentation(path: Path) -> None:
    prs = Presentation(str(path))
    for si, slide in enumerate(prs.slides):
        remove_empty_textboxes(slide)
        if si == 9:
            fix_systematic_ladder_slide(slide)
        title_idx = title_shape_index(slide)
        apply_text_colors(slide, title_idx)

    prs.save(str(path))


def main() -> None:
    if not PPTX_PATH.exists():
        raise SystemExit(f"Missing {PPTX_PATH}")
    polish_presentation(PPTX_PATH)
    print(f"Updated: {PPTX_PATH}")


if __name__ == "__main__":
    main()
